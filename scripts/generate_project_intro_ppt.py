from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent.parent
OUTPUT_DIR = BASE_DIR / "docs"
OUTPUT_PATH = OUTPUT_DIR / "RAG_LINEBOT_project_intro_5slides.pptx"
SETTINGS_IMG = BASE_DIR / "assets" / "settings.png"
KNOWLEDGE_IMG = BASE_DIR / "assets" / "knowledge.png"


NAVY = RGBColor(10, 25, 47)
NAVY_2 = RGBColor(17, 34, 64)
STEEL = RGBColor(0, 120, 212)
CYAN = RGBColor(41, 196, 255)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(243, 246, 251)
MUTED = RGBColor(180, 191, 211)
GREEN = RGBColor(18, 196, 138)
DARK_PANEL = RGBColor(21, 31, 56)
BORDER = RGBColor(49, 69, 104)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_full_background(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, title, accent_text=None):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = STEEL
    line.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(7.6), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = WHITE

    if accent_text:
        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(0.28), Inches(2.15), Inches(0.42)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY_2
        badge.line.color.rgb = BORDER
        tf = badge.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = accent_text
        r.font.name = "Segoe UI"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = CYAN


def add_footer(slide, text="RAG LINE Bot V2 | Project Overview"):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(4.5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED


def add_text(slide, left, top, width, height, text, *, size=24, color=WHITE, bold=False,
             font_name="Segoe UI", align=PP_ALIGN.LEFT, level=0, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    paragraphs = text.split("\n")
    first = True
    for line in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.level = level
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullet_list(slide, left, top, width, bullet_lines, *, size=23, color=OFF_WHITE, bullet_color=CYAN):
    tb = slide.shapes.add_textbox(left, top, width, Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, line in enumerate(bullet_lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return tb


def add_panel(slide, left, top, width, height, title, body_lines, accent=STEEL):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = DARK_PANEL
    panel.line.color.rgb = BORDER

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_text(slide, left + Inches(0.28), top + Inches(0.18), width - Inches(0.45), Inches(0.45), title,
             size=22, color=WHITE, bold=True)
    add_bullet_list(slide, left + Inches(0.32), top + Inches(0.68), width - Inches(0.55), body_lines,
                    size=17, color=OFF_WHITE)


def add_metric_chip(slide, left, top, width, text, value_color=CYAN):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.56))
    chip.fill.solid()
    chip.fill.fore_color.rgb = NAVY_2
    chip.line.color.rgb = BORDER
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = value_color


def add_image_card(slide, left, top, width, height, image_path, caption):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = DARK_PANEL
    frame.line.color.rgb = BORDER

    slide.shapes.add_picture(str(image_path), left + Inches(0.15), top + Inches(0.15), width=width - Inches(0.3), height=height - Inches(0.7))
    add_text(slide, left + Inches(0.18), top + height - Inches(0.42), width - Inches(0.36), Inches(0.25), caption,
             size=13, color=MUTED, bold=False)


# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_background(slide, NAVY)
hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.85), Inches(0.0), Inches(4.48), Inches(7.5))
hero.fill.solid()
hero.fill.fore_color.rgb = NAVY_2
hero.line.fill.background()

add_text(slide, Inches(0.72), Inches(1.1), Inches(7.5), Inches(1.2), "RAG LINE Bot V2",
         size=30, color=CYAN, bold=True)
add_text(slide, Inches(0.72), Inches(1.75), Inches(7.2), Inches(1.6), "智慧知識庫 Agent 系統",
         size=33, color=WHITE, bold=True)
add_text(slide, Inches(0.74), Inches(2.78), Inches(6.7), Inches(1.2),
         "以 FastAPI、Hybrid RAG 與模組化技能為核心，打造可決策、可維運、可擴充的 LINE 智慧助手平台。",
         size=19, color=OFF_WHITE)

add_metric_chip(slide, Inches(0.74), Inches(4.18), Inches(1.95), "FastAPI")
add_metric_chip(slide, Inches(2.88), Inches(4.18), Inches(2.1), "Hybrid RAG")
add_metric_chip(slide, Inches(5.16), Inches(4.18), Inches(2.18), "Skill System")

if SETTINGS_IMG.exists():
    slide.shapes.add_picture(str(SETTINGS_IMG), Inches(9.3), Inches(1.2), width=Inches(3.35), height=Inches(4.8))
add_text(slide, Inches(0.74), Inches(6.2), Inches(6.5), Inches(0.6),
         "專案簡報｜5 Slides Overview", size=14, color=MUTED, bold=False)
add_footer(slide)

# Slide 2: Positioning
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_background(slide, NAVY)
add_top_bar(slide, "專案定位", "Conversational AI + Knowledge Ops")
add_text(slide, Inches(0.72), Inches(1.0), Inches(5.4), Inches(0.7), "從聊天機器人進化成可決策的知識型 Agent",
         size=28, color=WHITE, bold=True)
add_text(slide, Inches(0.74), Inches(1.75), Inches(5.0), Inches(1.2),
         "這個專案不只是回覆訊息，而是讓系統能理解意圖、選擇技能、決定是否檢索，再自然地完成作答。",
         size=18, color=OFF_WHITE)

add_panel(slide, Inches(0.74), Inches(3.0), Inches(3.7), Inches(2.65), "要解決的問題", [
    "傳統 Bot 缺少推理與路由能力",
    "知識問答容易命中不足或答案生硬",
    "導入後常欠缺監控與可調校後台",
], accent=RGBColor(255, 122, 89))

add_panel(slide, Inches(4.82), Inches(3.0), Inches(3.7), Inches(2.65), "本專案的做法", [
    "以多階段 Agent 流程執行思考與動作",
    "用 Hybrid RAG 強化知識命中率與相關性",
    "以技能系統與控制台提高可擴充性",
], accent=STEEL)

add_panel(slide, Inches(8.9), Inches(3.0), Inches(3.7), Inches(2.65), "適用場景", [
    "FAQ / 文件查詢 / 內部知識助理",
    "LINE 客服自動化與營運支援",
    "需要繁中自然回覆的企業場景",
], accent=GREEN)
add_footer(slide)

# Slide 3: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_background(slide, NAVY)
add_top_bar(slide, "核心技術架構", "3 Pillars of the System")
add_text(slide, Inches(0.72), Inches(0.98), Inches(4.5), Inches(0.6), "三個關鍵能力，支撐高品質問答流程",
         size=27, color=WHITE, bold=True)

add_panel(slide, Inches(0.74), Inches(1.85), Inches(3.9), Inches(3.9), "1. Agent Router", [
    "根據意圖判斷是否需要檢索",
    "支援 READ_SKILL / PREPROCESS_QUERY / CALL_RAG",
    "以結構化 JSON 決策，提升流程穩定性",
], accent=CYAN)

add_panel(slide, Inches(4.72), Inches(1.85), Inches(3.9), Inches(3.9), "2. Hybrid RAG", [
    "Dense Retrieval + BM25 雙軌檢索",
    "RRF 融合不同來源結果",
    "Rerank 二次排序，讓相關內容更靠前",
], accent=STEEL)

add_panel(slide, Inches(8.7), Inches(1.85), Inches(3.9), Inches(3.9), "3. Skill System", [
    "每個技能擁有獨立 SKILL.md 規範",
    "可熱插拔啟用、停用或強制載入",
    "技能知識隔離，按需讀取參考內容",
], accent=GREEN)

flow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(2.95), Inches(5.95), Inches(7.3), Inches(0.58))
flow.fill.solid()
flow.fill.fore_color.rgb = NAVY_2
flow.line.color.rgb = BORDER
flow.text_frame.text = "User Query  →  Agent Decision  →  Retrieval / Skills  →  Natural Answer"
flow.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
flow.text_frame.paragraphs[0].font.size = Pt(17)
flow.text_frame.paragraphs[0].font.bold = True
flow.text_frame.paragraphs[0].font.color.rgb = OFF_WHITE
add_footer(slide)

# Slide 4: Product Experience
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_background(slide, NAVY)
add_top_bar(slide, "產品與管理體驗", "Admin Console & Operations")
add_text(slide, Inches(0.72), Inches(0.98), Inches(5.0), Inches(0.7), "不只回答問題，也讓系統可被管理",
         size=28, color=WHITE, bold=True)
add_text(slide, Inches(0.74), Inches(1.62), Inches(4.8), Inches(1.25),
         "前端控制台提供參數調整、技能管理、請求監控與知識庫維護，讓開發與營運能共同使用同一套工作台。",
         size=18, color=OFF_WHITE)

if SETTINGS_IMG.exists():
    add_image_card(slide, Inches(5.15), Inches(1.15), Inches(3.6), Inches(4.95), SETTINGS_IMG, "參數設定介面｜可調整 LINE、LLM、RAG 與服務連線")
if KNOWLEDGE_IMG.exists():
    add_image_card(slide, Inches(8.95), Inches(1.15), Inches(3.6), Inches(4.95), KNOWLEDGE_IMG, "知識庫管理介面｜支援文件維護、切片與向量化流程")

add_panel(slide, Inches(0.74), Inches(3.15), Inches(3.95), Inches(2.45), "營運面優勢", [
    "即時觀察 Agent 決策與檢索狀態",
    "快速調整 Top-K、溫度與技能開關",
    "降低上線後維護與排障成本",
], accent=CYAN)
add_footer(slide)

# Slide 5: Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_background(slide, NAVY)
add_top_bar(slide, "專案價值總結", "Ready for Extension and Deployment")
add_text(slide, Inches(0.72), Inches(1.0), Inches(6.4), Inches(0.7), "一套可擴充、可調校、可上線的智慧助手基礎建設",
         size=29, color=WHITE, bold=True)

summary_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.95), Inches(12.0), Inches(3.9))
summary_box.fill.solid()
summary_box.fill.fore_color.rgb = DARK_PANEL
summary_box.line.color.rgb = BORDER

add_text(slide, Inches(1.12), Inches(2.3), Inches(2.8), Inches(0.5), "更準", size=30, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.22), Inches(2.3), Inches(2.8), Inches(0.5), "更聰明", size=30, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(7.32), Inches(2.3), Inches(2.8), Inches(0.5), "更好維運", size=30, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(10.42), Inches(2.3), Inches(1.0), Inches(0.5), "", size=1)

add_text(slide, Inches(1.05), Inches(2.9), Inches(2.9), Inches(1.5),
         "Hybrid RAG、RRF 與 Rerank\n提升知識命中率與答案相關性",
         size=18, color=OFF_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(4.15), Inches(2.9), Inches(2.9), Inches(1.5),
         "Agent 決策與技能機制\n讓系統更接近任務導向助理",
         size=18, color=OFF_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(7.25), Inches(2.9), Inches(2.9), Inches(1.5),
         "後台設定、技能開關與監控能力\n降低導入與維護成本",
         size=18, color=OFF_WHITE, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.76), Inches(6.15), Inches(11.8), Inches(0.6),
         "結論：本專案已具備技術深度與產品化雛形，適合作為企業知識問答、LINE 智慧客服與可擴展 Agent 平台的基礎。",
         size=19, color=WHITE, bold=False)
add_footer(slide, "RAG LINE Bot V2 | Final Takeaways")

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT_PATH))
print(f"Generated: {OUTPUT_PATH}")
